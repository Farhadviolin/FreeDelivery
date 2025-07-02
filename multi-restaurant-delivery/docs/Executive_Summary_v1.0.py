from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Projektabschluss Delivery System v1.0"
slide.placeholders[1].text = "Alle Hauptfeatures, Security, Monitoring, Self-Healing, SLA, Accessibility, Doku & Go-Live bereit.\n\nSiehe Release Notes & Doku für Details.\n\nDanke an alle Teams!"
prs.save("Executive_Summary_v1.0.pptx")
